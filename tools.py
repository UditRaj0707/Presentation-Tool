# from langchain_core.tools import tool
from pptx import Presentation
from pptx.chart.data import XyChartData
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION, XL_LEGEND_POSITION
from typing import Optional
import matplotlib.pyplot as plt
import numpy as np
from abc import ABC
# import Annotated
from typing import Annotated
from pydantic import BaseModel
from langchain.tools import tool
from langchain.tools import StructuredTool
from tool_descriptions import TOOL_DESCRIPTIONS

import os
from load_dotenv import load_dotenv
load_dotenv()

SAVE_PATH = "output/presentation_test.pptx"
INPUT_PATH = None

class Presentationtools:
    """
    Presentationtools class
    """
    name: str = "presentation_tools"
    description: str = "Tools for interacting with PowerPoint slides"

    def __init__(self, mode="normal"):
        self.config = {
            "title_font_size": 40,
            "slide_width": 13.33,
            "slide_height": 7.5
        }
        self.mode = mode

    @staticmethod
    def use_presentation(folder_path: str, file_path: str):
        """
        Load an existing presentation file
        Args:
            file_path: Path to the presentation file
        """
        global INPUT_PATH, SAVE_PATH
        INPUT_PATH = folder_path + "/" + file_path
        SAVE_PATH = "output/" + file_path
        print(f"Presentation loaded from: {INPUT_PATH}")
        print(f"Presentation will be saved to: {SAVE_PATH}")

    @staticmethod
    def get_presentation():
        if INPUT_PATH is not None:
            return Presentation(INPUT_PATH)
        else:
            return Presentation()


    def add_image_slide(self, image_path: str, caption: str, title: str, insert_at: str = None):
        """
        Add a slide having image with caption to the presentation
        Args:
            image_path: Path to image file
            caption: caption text
            title: Title of the slide
        """
        if not hasattr(self, "prs"):
            self.prs = self.get_presentation()

        self.prs.slide_width = Inches(13.33) 
        self.prs.slide_height = Inches(7.5)

        if insert_at is None:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        else:
            slide = self.insert_slide(self.prs, layout_index=5, position=int(insert_at))

        title_shape = slide.shapes.title
        title_shape.text = title  
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.size = Pt(self.config["title_font_size"])
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = 2  # Center align (0=left, 1=center, 2=right)
        
        # Center the image horizontally by calculating left position
        img_width = 4  # inches
        left = (13.33 - img_width) / 2  # Center horizontally
        top = Inches(2)   
        width = 4
        height = 4
        pic = slide.shapes.add_picture(
            image_path,
            Inches(left),
            top,
            width=Inches(width),
            height=Inches(height)
        )
        
        if caption:
            # Center the caption box under the image
            left = (13.33 - img_width) / 2  # Match image position
            top = top + Inches(height) + Inches(0.5) 
            width = Inches(4)
            height = Inches(1)
            txBox = slide.shapes.add_textbox(Inches(left), top, width, height)
            tf = txBox.text_frame
            tf.text = caption
            tf.fit_text(font_family="Arial", max_size=12, italic=True)
            
        return self.save_presentation()
    
    def add_text_with_image_slide(self, text_content: str, image_path: str, title: str, insert_at: str = None, **kwargs):
        """
        Add a slide with text on the left half and an image on the right half.

        Args:
            text_content: Text content for the left side, paragraphs separated by double line breaks.
                            Example: "Paragraph 1.\n\nParagraph 2.\n\nParagraph 3."
            image_path: Path to the image file.
            title: Title of the slide.
        """
        if not hasattr(self, "prs"):
            self.prs = self.get_presentation()

        self.prs.slide_width = Inches(self.config['slide_width']) 
        self.prs.slide_height = Inches(self.config['slide_height'])

        if insert_at is None:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        else:
            slide = self.insert_slide(self.prs, layout_index=5, position=int(insert_at))

        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(self.config["title_font_size"])
        title_shape.text_frame.paragraphs[0].font.bold = True

        left_text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(4.5))
        text_frame = left_text_box.text_frame
        text_frame.word_wrap = True

        paragraphs = text_content.split('\n\n')
        for i, paragraph in enumerate(paragraphs):
            if i == 0:
                p = text_frame.paragraphs[0] 
            else:
                p = text_frame.add_paragraph()
            p.text = paragraph.strip()
            p.font.size = Pt(14)
            p.space_after = Pt(10) 

        img_left = Inches(5)  
        img_top = Inches(2)
        img_width = Inches(4) 
        img_height = Inches(4.5)
        slide.shapes.add_picture(image_path, img_left, img_top, width=img_width, height=img_height)

        # Save the presentation
        file_path = self.save_presentation()
        return f"Slide with image and text created and saved at: {file_path}"
    
    def add_bullet_slide(self, title: str, content: str, insert_at: str = None):
        """
        Create a slide with bullet points
        Args:
            title: Title of the slide 
            content: String of bullet points separated by semicolons
            Example: "First point; Second point; Third point"
        """
        if not hasattr(self, "prs"):
            self.prs = self.get_presentation()

        self.prs.slide_width = Inches(self.config['slide_width']) 
        self.prs.slide_height = Inches(self.config['slide_height'])
            
        if insert_at is None:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Use blank layout
        else:
            slide = self.insert_slide(self.prs, layout_index=5, position=int(insert_at))

        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title 
        title_shape.text_frame.paragraphs[0].font.size = Pt(self.config["title_font_size"])
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].alignment = 1
        
        # Create textbox for bullet points
        left = Inches(1)
        top = Inches(2) 
        width = Inches(10)
        height = Inches(5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = txBox.text_frame
        text_frame.word_wrap = True
        
        # Add bullet points
        bullet_points = [point.strip() for point in content.split(';')]
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = "• " + point
            p.font.size = Pt(20)
            # p.font.name = 'Arial'
            p.space_after = Pt(12)
        
        return self.save_presentation()
    
    def add_two_content_bullet_slide(self, title: str, left_content: str, right_content: str, insert_at: str = None):
        """
        Create a slide with two columns of bullet points
        Args:
            title: Title of the slide
            left_content: String with bullet points separated by semicolons
            right_content: String with bullet points separated by semicolons
            Example: 
                left_content = "Point 1; Point 2; Point 3"
                right_content = "Item A; Item B; Item C"
        """
        if not hasattr(self, "prs"):
            self.prs = self.get_presentation()
        
        # Set slide dimensions for 16:9 aspect ratio
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)

        if insert_at is None:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Using blank layout
        else:
            slide = self.insert_slide(self.prs, layout_index=5, position=int(insert_at))
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title  
        title_shape.text_frame.paragraphs[0].font.size = Pt(self.config["title_font_size"])
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].alignment = 1
        
        # Create left textbox
        left = Inches(1)
        top = Inches(2)
        width = Inches(5)
        height = Inches(4)
        left_box = slide.shapes.add_textbox(left, top, width, height)
        left_frame = left_box.text_frame
        left_frame.word_wrap = True

        # Add left content with bullet points
        for i, point in enumerate(left_content.split(';')):
            if i == 0:
                p = left_frame.paragraphs[0]
            else:
                p = left_frame.add_paragraph()
            p.text = "• " + point.strip()  # Add bullet character
            p.font.size = Pt(20)
            p.space_before = Pt(12)
        
        # Create right textbox
        left = Inches(7)  # Position this on the right side
        right_box = slide.shapes.add_textbox(left, top, width, height)
        right_frame = right_box.text_frame
        right_frame.word_wrap = True

        # Add right content with bullet points
        for i, point in enumerate(right_content.split(';')):
            if i == 0:
                p = right_frame.paragraphs[0]
            else:
                p = right_frame.add_paragraph()
            p.text = "• " + point.strip()  # Add bullet character
            p.font.size = Pt(20)
            p.space_before = Pt(12)
        
        return self.save_presentation()
        
    def add_table_slide(self, table_data: str, title: str, insert_at: str = None):
        """
        Add a comparison table to the presentation from string input
        Args:
            table_data: String in format "header1, header2; value1, value2; value3, value4"
            title: Title of the slide
            Example: "Car, Bike; BMW, Harley; Audi, Ducati; Mercedes, Honda"
        """
        if not hasattr(self, "prs"):
            self.prs = self.get_presentation() 

        self.prs.slide_width = Inches(self.config['slide_width']) 
        self.prs.slide_height = Inches(self.config['slide_height'])

        rows = table_data.split(';')
        headers = rows[0].split(',')
        values = [row.split(',') for row in rows[1:]]
        
        table_data = {}
        for i, header in enumerate(headers):
            table_data[header] = [row[i] for row in values]
        
        if insert_at is None:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        else:
            slide = self.insert_slide(self.prs, layout_index=5, position=int(insert_at))

        title_shape = slide.shapes.title
        title_shape.text = title 
        title_shape.text_frame.paragraphs[0].font.size = Pt(self.config["title_font_size"])
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].alignment = 1

        # Calculate dimensions and position
        rows = len(list(table_data.values())[0]) + 1
        cols = len(table_data.keys())
        
        # Calculate table width and height based on content
        max_text_len = max(
            max(len(str(val)) for vals in table_data.values() for val in vals),
            max(len(str(key)) for key in table_data.keys())
        )
        
        # Base cell dimensions
        cell_width = max(1.5, max_text_len * 0.15)  # Minimum 1.5 inches
        table_width = cell_width * cols
        cell_height = 0.4  # Base height per cell
        table_height = cell_height * rows
        
        # Maximum allowed dimensions
        max_width = 11  # Maximum table width (leaving margins)
        max_height = 5  # Maximum table height (leaving space for title)
        
        # Only scale if table exceeds slide bounds
        if table_width > max_width:
            scale = max_width / table_width
            table_width = max_width
            cell_width *= scale
            
        if table_height > max_height:
            scale = max_height / table_height
            table_height = max_height
            cell_height *= scale
        
        # Center the table on slide
        left = (13.33 - table_width) / 2  # Center horizontally
        top = (7.5 - table_height) / 2 + 0.5  # Center vertically with offset for title

        table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), 
                                     Inches(table_width), Inches(table_height)).table
        
        # Style the table
        for i, key in enumerate(table_data.keys()):
            cell = table.cell(0, i)
            cell.text = key.strip()
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)
            
            for j, value in enumerate(table_data[key]):
                cell = table.cell(j + 1, i)
                cell.text = value.strip()
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(11)
        
        # Apply cell padding and alignment
        for row in table.rows:
            for cell in row.cells:
                cell.margin_left = cell.margin_right = Inches(0.1)
                cell.margin_top = cell.margin_bottom = Inches(0.05)
                cell.vertical_anchor = 1  # Center vertical alignment
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = 1  # Center horizontal alignment
        
        file_path = self.save_presentation()
        return file_path


    def add_bar_chart(self, categories_str: str, series_data_str: str, title: str, insert_at: str = None) -> str:
        """
        Create a slide with a clustered bar chart for comparing multiple data series.

        Args:
            categories_str: String of categories separated by commas.
                            Example: "East, West, Midwest"
            series_data_str: String with series data in the format:
                            "Series1: 19.2, 21.4, 16.7; Series2: 22.3, 28.6, 15.2; Series3: 20.4, 26.3, 14.2"
            title: Title of the slide.

        Example:
            categories_str = "East, West, Midwest"
            series_data_str = "Q1: 19.2, 21.4, 16.7; Q2: 22.3, 28.6, 15.2; Q3: 20.4, 26.3, 14.2"
        """
        if not hasattr(self, "prs"):
            self.prs = self.get_presentation()

        self.prs.slide_width = Inches(self.config['slide_width']) 
        self.prs.slide_height = Inches(self.config['slide_height'])
        # Create a new slide
        if insert_at is None:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        else:
            slide = self.insert_slide(self.prs, layout_index=5, position=int(insert_at))

        # Set slide title with custom font size
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(24)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].alignment = 1

        # Initialize the chart data
        chart_data = ChartData()
        chart_data.categories = [c.strip() for c in categories_str.split(",")]

        # Process multiple series data
        series_entries = series_data_str.split(";")
        for series_entry in series_entries:
            series_name, series_values = series_entry.split(":")
            values = [float(v.strip()) for v in series_values.split(",")]
            chart_data.add_series(series_name.strip(), values)

        # Add the clustered bar chart
        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False

        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels

        data_labels.font.size = Pt(6)
        data_labels.font.color.rgb = RGBColor(0, 0, 0)
        data_labels.position = XL_LABEL_POSITION.INSIDE_END


        # Save the presentation
        file_path = self.save_presentation()
        return f"Slide with clustered bar chart created and saved at: {file_path}"



    def add_line_chart(self, categories_str: str, series_data_str: str, title: str, insert_at: str = None, **kwargs) -> str:
        """
        Create a slide with a multi-series line chart.

        Args:
            categories_str: String of categories separated by commas.
                            Example: "Q1 Sales, Q2 Sales, Q3 Sales"
            series_data_str: String of series data in the format:
                            "West: 30, 28, 35; East: 25, 30, 20; Midwest: 20, 18, 25"
            title: Title of the slide.
        """
        if not hasattr(self, "prs"):
            self.prs = self.get_presentation()

        self.prs.slide_width = Inches(self.config['slide_width']) 
        self.prs.slide_height = Inches(self.config['slide_height'])

        if insert_at is None:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        else:
            slide = self.insert_slide(self.prs, layout_index=5, position=int(insert_at))

        title_shape = slide.shapes.title
        title_shape.text = title 
        title_shape.text_frame.paragraphs[0].font.size = Pt(self.config["title_font_size"])
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].alignment = 1

        chart_data = CategoryChartData()
        chart_data.categories = [c.strip() for c in categories_str.split(",")]

        series_entries = series_data_str.split(";")
        for series_entry in series_entries:
            series_name, series_values = series_entry.split(":")
            values = [float(v.strip()) for v in series_values.split(",")]
            chart_data.add_series(series_name.strip(), values)

        x, y, cx, cy = Inches(3.5), Inches(2), Inches(6), Inches(4.5)
        slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data)

        file_path = self.save_presentation()
        print(f"Slide with multi-series line chart created and saved at: {file_path}")
        return f"Slide with multi-series line chart created and saved at: {file_path}"

    
    def add_pie_chart(self, categories_str: str, values_str: str, right_content: str, title: str, plot_name: str, insert_at: str = None) -> str:
        """
        Create a slide with a pie chart including well-formatted category labels.
        Args:
            categories_str: String of categories separated by commas
            values_str: String of values separated by commas
            right_content: String of bullet points separated by semicolons
            title: Title of the slide
            plot_name: Title of the pie chart
        """
        if not hasattr(self, "prs"):
            self.prs = self.get_presentation()

        self.prs.slide_width = Inches(self.config['slide_width']) 
        self.prs.slide_height = Inches(self.config['slide_height'])

        if insert_at is None:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        else:
            slide = self.insert_slide(self.prs, layout_index=5, position=int(insert_at))

        title_shape = slide.shapes.title
        title_shape.text = title  
        title_shape.text_frame.paragraphs[0].font.size = Pt(self.config["title_font_size"])
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].alignment = 1

        chart_data = CategoryChartData()
        chart_data.categories = [c.strip() for c in categories_str.split(",")]
        chart_data.add_series(plot_name, (float(v.strip()) for v in values_str.split(",")))

        x, y, cx, cy = Inches(1), Inches(2), Inches(6), Inches(4.5)
        chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data)
        chart = chart_shape.chart 

        plot = chart.plots[0]
        plot.has_data_labels = True  
        data_labels = plot.data_labels
        data_labels.show_category_name = True
        data_labels.show_percentage = True
        data_labels.separator = "\n" 
        data_labels.position = XL_DATA_LABEL_POSITION.BEST_FIT 
        
        # Create textbox for bullet points on right side
        right_textbox = slide.shapes.add_textbox(Inches(8), Inches(2), Inches(4), Inches(4.5))
        right_frame = right_textbox.text_frame
        right_frame.word_wrap = True  # Enable word wrap
        right_frame.clear()  # Clear existing text

        # Add bullet points
        right_points = [point.strip() for point in right_content.split(';')]
        for i, point in enumerate(right_points):
            if i == 0:
                # Set the first paragraph's text
                p = right_frame.paragraphs[0]
                p.text = f"{i+1}. " + point + "\n"
            else:
                # Add new paragraphs for subsequent bullet points
                p = right_frame.add_paragraph()
                p.text = f"{i+1}. " + point + "\n"
            
            # Format bullet points
            p.level = 0
            p.font.size = Pt(20)  # Set consistent font size

        file_path = self.save_presentation()
        return f"Slide with pie chart created and saved at: {file_path}"
    
    def add_area_chart(self, categories_str: str, values_str: str, title: str, plot_name : str, insert_at: str = None) -> str:
        """
        Create a slide with an area chart
        Args:
            categories_str: String of categories separated by commas
            values_str: String of values separated by commas
            title: Title of the slide
            plot_name: Title of the area chart
            Example: "A, B, C, D", "1, 2, 3, 4", "Area Chart Slide"
        """
        if not hasattr(self, "prs"):
            self.prs = self.get_presentation()

        self.prs.slide_width = Inches(self.config['slide_width']) 
        self.prs.slide_height = Inches(self.config['slide_height'])

        if insert_at is None:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        else:
            slide = self.insert_slide(self.prs, layout_index=5, position=int(insert_at))

        title_shape = slide.shapes.title
        title_shape.text = title  
        title_shape.text_frame.paragraphs[0].font.size = Pt(self.config["title_font_size"])
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].alignment = 1

        chart_data = CategoryChartData()
        chart_data.categories = [c.strip() for c in categories_str.split(",")]
        chart_data.add_series(plot_name, (float(v.strip()) for v in values_str.split(",")))

        x, y, cx, cy = Inches(4), Inches(2.5), Inches(6), Inches(4.5)
        slide.shapes.add_chart(XL_CHART_TYPE.AREA, x, y, cx, cy, chart_data)

        file_path = self.save_presentation()
        return f"Slide with area chart created and saved at: {file_path}"
    
    def add_scatter_chart(self, input_x: str, input_y: str, title: str, plot_title: str, insert_at: str = None) -> str:
        """
    Creates a scatter plot slide with labeled axes.

    Parameters:
    - input_x (str): X-axis label and values in the format "Label; 1, 2, 3, 4" or just "1, 2, 3, 4".
    - input_y (str): Y-axis label and values in the format "Label; 5, 6, 7, 8" or just "5, 6, 7, 8".
    - title (str): The title of the slide.
    - plot_title (str): The title of the scatter plot.
    
    Returns:
    - str: The file path of the saved presentation slide.

    Example:
    - add_scatter_chart("X-Axis; 1, 2, 3, 4", "Y-Axis; 5, 6, 7, 8", "Scatter Plot Slide", "Scatter Plot")
    - add_scatter_chart("1, 2, 3, 4", "5, 6, 7, 8", "Scatter Plot Slide", "Scatter Plot")
    """
        try:
            if ";" in input_x:
                x_label, x_values_str = input_x.split(";")
            else:
                x_label, x_values_str = "X-Axis", input_x
            
            if ";" in input_y:
                y_label, y_values_str = input_y.split(";")
            else:
                y_label, y_values_str = "Y-Axis", input_y
            
            x_values = [float(x.strip()) for x in x_values_str.split(",")]
            y_values = [float(y.strip()) for y in y_values_str.split(",")]
        except ValueError:
            return "Error: Inputs must be in the format 'Label; 1, 2, 3, 4' or '1, 2, 3, 4' with numeric values."
        
        if not hasattr(self, "prs"):
            self.prs = self.get_presentation()

        self.prs.slide_width = Inches(self.config['slide_width']) 
        self.prs.slide_height = Inches(self.config['slide_height'])

        if insert_at is None:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        else:
            slide = self.insert_slide(self.prs, layout_index=5, position=int(insert_at))

        title_shape = slide.shapes.title
        title_shape.text = title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.size = Pt(self.config["title_font_size"]) 
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = 1  

        chart_data = XyChartData()
        series = chart_data.add_series(plot_title)

        for x, y in zip(x_values, y_values):
            series.add_data_point(x, y)

        x, y, cx, cy = Inches(4), Inches(2.5), Inches(6), Inches(4.5)
        chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data)
        chart = chart_shape.chart

        category_axis = chart.category_axis
        value_axis = chart.value_axis
        category_axis.has_title = True
        value_axis.has_title = True
        category_axis.axis_title.text_frame.text = x_label.strip()
        value_axis.axis_title.text_frame.text = y_label.strip()

        file_path = self.save_presentation()
        return f"Slide with scatter chart (with axis labels) created and saved at: {file_path}"

    
    # def get_slide_layout(self):
    #     for idx, layout in enumerate(self.prs.slide_layouts):
    #         print(f"Layout {idx}: {layout.name}")
    #         slide = self.prs.slides.add_slide(layout)
    #         for placeholder in slide.placeholders:
    #             print(f"\tPlaceholder {placeholder.placeholder_format.idx}: {placeholder.name}")

    def add_waterfall_chart(self, categories_str: str, values_str: str, title: str, totals_str: Optional[str] = None, insert_at: str = None) -> None:
        """
        Create and display a waterfall chart.

        Args:
            categories_str: Comma-separated string of category names (e.g., "Category 1, Category 2, Category 3").
            values_str: Comma-separated string of incremental values (e.g., "100, 20, 50, -40, 130, -60, 70, 140").
            totals_str: Comma-separated string of total values (e.g., "100, None, None, None, 130, None, None, 140").
                        If not provided, totals will be calculated automatically.
        """
        # Parse inputs
        categories = [c.strip() for c in categories_str.split(",")]
        values = [float(v.strip()) for v in values_str.split(",")]
        if totals_str:
            totals = [float(v.strip()) if v.strip().lower() != "none" else None for v in totals_str.split(",")]
        else:
            # Auto-calculate totals
            totals = []
            cumulative_sum = values[0]
            for i, value in enumerate(values):
                if i == 0:
                    totals.append(value)
                elif value == cumulative_sum:
                    totals.append(value)
                    cumulative_sum = value
                else:
                    totals.append(None)
                    cumulative_sum += value

        plot_path = self.plot_waterfall_chart(categories, values, totals)

        if not hasattr(self, "prs"):
            self.prs = self.get_presentation()

        self.prs.slide_width = Inches(self.config['slide_width']) 
        self.prs.slide_height = Inches(self.config['slide_height'])   
        if insert_at is None:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        else:
            slide = self.insert_slide(self.prs, layout_index=5, position=int(insert_at))

        title_shape = slide.shapes.title
        title_shape.text = title  
        title_shape.text_frame.paragraphs[0].font.size = Pt(self.config["title_font_size"])
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].alignment = 1

        left = Inches(0.5)  
        top = Inches(1.5)   
        width = 12
        height = 5.5
        pic = slide.shapes.add_picture(
            plot_path,
            left,
            top,
            width=Inches(width),
            height=Inches(height)
        )

        file_path = self.save_presentation()
        return f"Slide with waterfall chart created and saved at: {file_path}"
    
    # add a function called add_title_slide in which you can add a title slide to the presentation, title should be in the center of the slide
    def add_title_slide(self, title: str, insert_at: str = None):
        """
        Add a title slide to the presentation with the given title centered on a 16:9 slide.
        Args:
            title: Title of the slide
        """
        if not hasattr(self, "prs"):
            self.prs = self.get_presentation()
            
        # Set slide dimensions for 16:9 aspect ratio 
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)

        if insert_at is None:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5]) # Use blank layout
        else:
            slide = self.insert_slide(self.prs, layout_index=5, position=int(insert_at))

        # Calculate center position
        left = Inches(1)
        top = Inches(2.5) # Center vertically
        width = Inches(11.33) # Full width minus margins
        height = Inches(2)

        # Add centered textbox
        txBox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = txBox.text_frame
        
        # Add and format text
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.alignment = 2 # Center alignment
        
        file_path = self.save_presentation()
        return f"Title slide created and saved at: {file_path}"


    def save_presentation(self):
        """
        Save the presentation to a file
        """
        if self.mode == "normal":
            print(f"self.save_path: {SAVE_PATH}")
            self.prs.save(SAVE_PATH)
            return SAVE_PATH
        else:
            save_path = f"output/{self.mode}.pptx"
            print(f"self.save_path: {save_path}")
            self.prs.save(save_path)
            return save_path
    
    def get_tools(self):
        print("🔍 Binding tools...")
        return [
            StructuredTool.from_function(
                func=lambda **kwargs: self.add_image_slide(**kwargs),
                name="add_image_slide",
                description=TOOL_DESCRIPTIONS.get("add_image_slide")
            ),
            StructuredTool.from_function(
                func=lambda **kwargs: self.add_text_with_image_slide(**kwargs),
                name="add_text_with_image_slide",
                description=TOOL_DESCRIPTIONS.get("add_text_with_image_slide")
            ),
            StructuredTool.from_function(
                func=lambda **kwargs: self.add_bullet_slide(**kwargs),
                name="add_bullet_slide",
                description=TOOL_DESCRIPTIONS.get("add_bullet_slide")
            ),
            StructuredTool.from_function(
                func=self.add_two_content_bullet_slide,
                name="add_two_content_bullet_slide",
                description=TOOL_DESCRIPTIONS.get("add_two_content_bullet_slide")
            ),
            StructuredTool.from_function(
                func=self.add_table_slide,
                name="add_table_slide",
                description=TOOL_DESCRIPTIONS.get("add_table_slide")
            ),
            StructuredTool.from_function(
                func=self.add_bar_chart,
                name="add_bar_chart",
                description=TOOL_DESCRIPTIONS.get("add_bar_chart")
            ),
            StructuredTool.from_function(
                func=self.add_line_chart,
                name="add_line_chart",
                description=TOOL_DESCRIPTIONS.get("add_line_chart")
            ),
            StructuredTool.from_function(
                func=self.add_pie_chart,
                name="add_pie_chart",
                description=TOOL_DESCRIPTIONS.get("add_pie_chart")
            ),
            StructuredTool.from_function(
                func=self.add_area_chart,
                name="add_area_chart", 
                description=TOOL_DESCRIPTIONS.get("add_area_chart")
            ),
            StructuredTool.from_function(
                func=self.add_scatter_chart,
                name="add_scatter_chart",
                description=TOOL_DESCRIPTIONS.get("add_scatter_chart")
            ),
            StructuredTool.from_function(
                func=self.add_waterfall_chart,
                name="add_waterfall_chart",
                description=TOOL_DESCRIPTIONS.get("add_waterfall_chart")
            ),
            StructuredTool.from_function(
                func=self.add_title_slide,
                name="add_title_slide",
                description="Add a title slide to the presentation"
            )
        ]

    @staticmethod
    def plot_waterfall_chart(categories, values, totals, filename="output/waterfall_chart.png"):
        labels = [f"{v:+}" for v in values]  # Add "+" or "-" to the value labels

        # Initialize cumulative values
        cumulative = [0]  # Start from baseline (0)
        bar_positions = []  # Track where each bar should start

        for i, value in enumerate(values):
            if totals[i] is not None:  # Total bar resets to 0
                cumulative.append(totals[i])
                bar_positions.append(0)
            else:  # Incremental values
                bar_positions.append(cumulative[-1])
                cumulative.append(cumulative[-1] + value)

        # Assign colors: green for totals, blue for increases, orange for decreases
        colors = []
        for i, value in enumerate(values):
            if totals[i] is not None:  # Total
                colors.append("green")
            elif value > 0:  # Increase
                colors.append("steelblue")
            else:  # Decrease
                colors.append("darkorange")

        # Plot
        fig, ax = plt.subplots(figsize=(16, 6))
        bars = ax.bar(categories, values, bottom=bar_positions, color=colors)

        # Add labels to each bar
        for bar, label in zip(bars, labels):
            height = bar.get_height()
            if height > 0:
                ax.text(bar.get_x() + bar.get_width() / 2, bar.get_y() + height - 10, label, ha="center", va="bottom", fontsize=10)
            else:
                ax.text(bar.get_x() + bar.get_width() / 2, bar.get_y() + height - 10, label, ha="center", va="top", fontsize=10)

        # Add title and legend
        ax.set_title("Waterfall Chart", fontsize=14)
        ax.set_ylabel("Values", fontsize=12)
        ax.legend(
            handles=[
                plt.Rectangle((0, 0), 1, 1, color="green", label="Total"),
                plt.Rectangle((0, 0), 1, 1, color="steelblue", label="Increase"),
                plt.Rectangle((0, 0), 1, 1, color="darkorange", label="Decrease"),
            ],
            loc="upper right",
        )

        # Add grid
        ax.grid(axis="y", linestyle="--", alpha=0.7)

        # Show plot
        plt.tight_layout()
        plt.savefig(filename)
        
        return filename
    
    @staticmethod
    def insert_slide(prs, layout_index, position):
        """
        Insert a slide at a specific position in the presentation.

        Args:
            prs (Presentation): The presentation object.
            layout_index (int): The index of the slide layout to use.
            position (int): The 0-based position where the slide should be inserted.

        Returns:
            slide: The newly inserted slide.
        """
        # Add a new slide (temporarily at the end)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        
        # Access the slide ID list
        slide_id_list = prs.slides._sldIdLst

        # Get the XML element for the new slide
        new_slide_id = slide_id_list[-1]
        
        # Remove the new slide from the slide ID list
        slide_id_list.remove(new_slide_id)

        # Insert the new slide at the desired position
        slide_id_list.insert(position, new_slide_id)

        return prs.slides[position]



if __name__ == "__main__":
    tool = Presentationtools()




    
    